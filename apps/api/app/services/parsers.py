from __future__ import annotations

import json
import re
import unicodedata
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup

from app.core.config import get_settings
from app.core.utils import source_type_from_path


@dataclass
class ParsedSection:
    title: str
    text: str
    page_number: int | None = None
    section: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


def detect_source_type(path: Path) -> str:
    return source_type_from_path(path)


CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
MOJIBAKE_MARKERS = (
    "\ufffd",
    "\u00c3",
    "\u00c2",
    "\u00e2",
    "\u9208",
    "\u9365",
    "\u9429",
    "\u95b3",
    "\u951f",
    "\u7d34",
    "\u6d93",
    "\u934f",
    "\u704f",
    "\u93c4",
    "\u9a9e",
)
MOJIBAKE_MARKER_RE = re.compile("|".join(re.escape(marker) for marker in MOJIBAKE_MARKERS))
LATIN_WORD_HYPHEN_BREAK_RE = re.compile(r"(?<=[A-Za-z])-\n(?=[a-z])")
SOFT_SINGLE_NEWLINE_RE = re.compile(r"(?<![\n.!?:;。！？：；])\n(?!\s*(?:\n|[#>*+\-]|\d+[.)]))")
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "big5", "cp1252", "latin-1")


def _safe_import_ftfy():
    try:
        import ftfy

        return ftfy
    except Exception:
        return None


def _detect_encoding(raw: bytes) -> tuple[str | None, float | None]:
    try:
        from charset_normalizer import from_bytes

        match = from_bytes(raw).best()
        if match is not None and match.encoding:
            coherence = getattr(match, "percent_coherence", None)
            return str(match.encoding), float(coherence) if coherence is not None else None
    except Exception:
        pass
    return None, None


def _mojibake_score(text: str) -> float:
    if not text:
        return 0.0
    markers = len(MOJIBAKE_MARKER_RE.findall(text))
    replacement = text.count("\ufffd")
    controls = len(CONTROL_CHAR_RE.findall(text))
    return (markers * 2.0 + replacement * 3.0 + controls * 2.0) / max(len(text), 1)


def _repair_mojibake_candidate(text: str) -> tuple[str, bool]:
    original_score = _mojibake_score(text)
    if original_score <= 0:
        return text, False

    candidates = [text]
    ftfy = _safe_import_ftfy()
    if ftfy is not None:
        try:
            candidates.append(ftfy.fix_text(text, normalization="NFC"))
        except Exception:
            pass

    # Typical PDF / web extraction failure: UTF-8 bytes were decoded as a
    # legacy code page, producing CJK-looking garbage such as "閺嶇绺?.
    for encoding in ("gb18030", "big5", "cp1252", "latin-1"):
        try:
            candidates.append(text.encode(encoding, errors="strict").decode("utf-8", errors="strict"))
        except Exception:
            continue

    best = min(candidates, key=_mojibake_score)
    best_score = _mojibake_score(best)
    if best != text and best_score + 0.002 < original_score:
        return best, True
    return text, False


def decode_text_bytes(raw: bytes) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {}
    detected, coherence = _detect_encoding(raw)
    if detected:
        metadata["encoding_detected"] = detected
    if coherence is not None:
        metadata["encoding_coherence"] = round(coherence, 3)

    preferred = []
    detected_usable = detected and (coherence is None or coherence >= 20.0 or detected.lower().replace("_", "-") in TEXT_ENCODINGS)
    if detected_usable:
        preferred.append(detected)
    preferred.extend(encoding for encoding in TEXT_ENCODINGS if encoding.lower() not in {item.lower() for item in preferred})

    last_error: Exception | None = None
    for encoding in preferred:
        try:
            text = raw.decode(encoding)
            metadata.setdefault("encoding_used", encoding)
            return text, metadata
        except (LookupError, UnicodeDecodeError) as exc:
            last_error = exc
            continue
    metadata["encoding_used"] = "utf-8-ignore"
    if last_error is not None:
        metadata["encoding_error"] = str(last_error)
    return raw.decode("utf-8", errors="ignore"), metadata


def clean_extracted_text(text: str, *, source_type: str | None = None) -> tuple[str, dict[str, Any]]:
    flags: list[str] = []
    original = text or ""
    cleaned = original.replace("\r\n", "\n").replace("\r", "\n").replace("\ufeff", "")
    if cleaned != original:
        flags.append("normalized_line_endings_or_bom")
    before_controls = cleaned
    cleaned = CONTROL_CHAR_RE.sub("", cleaned)
    if cleaned != before_controls:
        flags.append("removed_control_chars")
    cleaned = unicodedata.normalize("NFC", cleaned)

    repaired, repaired_mojibake = _repair_mojibake_candidate(cleaned)
    if repaired_mojibake:
        cleaned = repaired
        flags.append("mojibake_repaired")

    if source_type in {"pdf", "image", "ocr"}:
        before_layout = cleaned
        cleaned = LATIN_WORD_HYPHEN_BREAK_RE.sub("", cleaned)
        cleaned = SOFT_SINGLE_NEWLINE_RE.sub(" ", cleaned)
        if cleaned != before_layout:
            flags.append("normalized_pdf_ocr_linebreaks")

    before_space = cleaned
    cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    cleaned = re.sub(r"\n{4,}", "\n\n\n", cleaned)
    cleaned = cleaned.strip()
    if cleaned != before_space.strip():
        flags.append("normalized_whitespace")

    metadata = {
        "text_cleaning_flags": sorted(set(flags)),
        "mojibake_repaired": repaired_mojibake,
        "mojibake_score_before": round(_mojibake_score(original), 6),
        "mojibake_score_after": round(_mojibake_score(cleaned), 6),
    }
    return cleaned, metadata


def load_text_with_metadata(path: Path) -> tuple[str, dict[str, Any]]:
    decoded, metadata = decode_text_bytes(path.read_bytes())
    cleaned, cleaning = clean_extracted_text(decoded, source_type=detect_source_type(path))
    return cleaned, {**metadata, **cleaning}


def load_text(path: Path) -> str:
    return load_text_with_metadata(path)[0]


def _clean_section(section: ParsedSection, source_type: str, common_metadata: dict[str, Any] | None = None) -> ParsedSection:
    text, cleaning = clean_extracted_text(section.text, source_type=source_type)
    title, title_cleaning = clean_extracted_text(section.title, source_type=source_type)
    section_label = section.section
    if section_label is not None:
        section_label, _ = clean_extracted_text(section_label, source_type=source_type)
    flags = sorted(set((cleaning.get("text_cleaning_flags") or []) + (title_cleaning.get("text_cleaning_flags") or [])))
    metadata = {
        **(common_metadata or {}),
        **section.metadata,
        **cleaning,
        "text_cleaning_flags": flags,
        "mojibake_repaired": bool(cleaning.get("mojibake_repaired") or title_cleaning.get("mojibake_repaired")),
    }
    return ParsedSection(
        title=title,
        text=text,
        page_number=section.page_number,
        section=section_label,
        metadata=metadata,
    )


def clean_parsed_sections(sections: list[ParsedSection], source_type: str, common_metadata: dict[str, Any] | None = None) -> list[ParsedSection]:
    return [
        cleaned
        for section in sections
        for cleaned in [_clean_section(section, source_type, common_metadata)]
        if cleaned.text
    ]


def parse_markdown(path: Path) -> list[ParsedSection]:
    text, text_metadata = load_text_with_metadata(path)
    lines = text.splitlines()
    sections: list[ParsedSection] = []
    current_title = path.stem
    buffer: list[str] = []
    for line in lines:
        if line.strip().startswith("#"):
            if buffer:
                sections.append(
                    ParsedSection(
                        title=current_title,
                        text="\n".join(buffer).strip(),
                        section=current_title,
                        metadata={"content_kind": "markdown", **text_metadata},
                    )
                )
                buffer = []
            current_title = line.lstrip("#").strip() or current_title
            continue
        buffer.append(line)
    if buffer:
        sections.append(
            ParsedSection(
                title=current_title,
                text="\n".join(buffer).strip(),
                section=current_title,
                metadata={"content_kind": "markdown", **text_metadata},
            )
        )
    return [section for section in sections if section.text]


def parse_text(path: Path) -> list[ParsedSection]:
    text, text_metadata = load_text_with_metadata(path)
    return [ParsedSection(title=path.stem, text=text.strip(), section=path.stem, metadata={"content_kind": "text", **text_metadata})]


def parse_html(path: Path) -> list[ParsedSection]:
    html, text_metadata = load_text_with_metadata(path)
    soup = BeautifulSoup(html, "html.parser")
    title = soup.title.text.strip() if soup.title and soup.title.text else path.stem
    text = soup.get_text("\n", strip=True)
    return [ParsedSection(title=title, text=text, section=title, metadata={"content_kind": "html", **text_metadata})]


def _detect_formula(text: str) -> bool:
    """Detect formula-heavy text without rewriting math symbols."""
    formula_chars = set("∑∫∂√∞≈≠≤≥±×÷∈∉⊂⊆∪∩→←↔∀∃∇αβγδθηλμπρστυφχψωΓΔΘΛΠΣΦΨΩ")
    if not text:
        return False
    ratio = sum(1 for c in text if c in formula_chars) / len(text)
    return ratio > 0.03


def parse_pdf(path: Path) -> list[ParsedSection]:
    import fitz

    sections: list[ParsedSection] = []
    with fitz.open(path) as document:
        for idx, page in enumerate(document, start=1):
            try:
                text = page.get_text("markdown").strip()
            except Exception:
                text = page.get_text("text").strip()

            if not text:
                continue

            # Preserve table/formula hints for downstream chunk metadata.
            lines = text.splitlines()
            has_table = any("|" in line and "---" in line for line in lines)
            has_formula = _detect_formula(text)

            page_title = lines[0][:120] if lines else ""
            metadata: dict[str, Any] = {"content_kind": "pdf_page"}
            if has_table:
                metadata["has_table"] = True
            if has_formula:
                metadata["has_formula"] = True

            sections.append(
                ParsedSection(
                    title=page_title or f"{path.stem} p.{idx}",
                    text=text,
                    page_number=idx,
                    section=page_title or path.stem,
                    metadata=metadata,
                )
            )
    return sections


def parse_presentation(path: Path) -> list[ParsedSection]:
    from pptx import Presentation

    presentation = Presentation(path)
    sections: list[ParsedSection] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                fragments.append(shape.text.strip())
        text = "\n".join(fragment for fragment in fragments if fragment)
        if text:
            lines = text.splitlines()
            sections.append(
                ParsedSection(
                    title=lines[0][:120] if lines else f"{path.stem} slide {idx}",
                    text=text,
                    page_number=idx,
                    section=f"slide-{idx}",
                    metadata={"content_kind": "slide"},
                )
            )
    return sections


def parse_docx(path: Path) -> list[ParsedSection]:
    from docx import Document

    document = Document(path)
    sections: list[ParsedSection] = []
    current_title = path.stem
    buffer: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and "Heading" in paragraph.style.name:
            if buffer:
                sections.append(
                    ParsedSection(
                        title=current_title,
                        text="\n".join(buffer),
                        section=current_title,
                        metadata={"content_kind": "doc_section"},
                    )
                )
                buffer = []
            current_title = text
            continue
        buffer.append(text)
    if buffer:
        sections.append(
            ParsedSection(
                title=current_title,
                text="\n".join(buffer),
                section=current_title,
                metadata={"content_kind": "doc_section"},
            )
        )
    return sections


def parse_notebook(path: Path) -> list[ParsedSection]:
    decoded, text_metadata = decode_text_bytes(path.read_bytes())
    notebook = json.loads(decoded)
    sections: list[ParsedSection] = []
    cell_buffer: list[str] = []
    current_title = path.stem
    current_index = 0
    for cell in notebook.get("cells", []):
        current_index += 1
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue
        if cell.get("cell_type") == "markdown":
            heading = next((line.lstrip("# ").strip() for line in source.splitlines() if line.strip().startswith("#")), None)
            if heading:
                if cell_buffer:
                    sections.append(
                        ParsedSection(
                            title=current_title,
                            text="\n".join(cell_buffer),
                            section=current_title,
                            metadata={"cell_index": current_index, "content_kind": "markdown", **text_metadata},
                        )
                    )
                    cell_buffer = []
                current_title = heading
            cell_buffer.append(source.strip())
        else:
            outputs: list[str] = []
            for output in cell.get("outputs", []):
                if "text" in output:
                    outputs.append("".join(output["text"]))
                elif "data" in output and "text/plain" in output["data"]:
                    outputs.append("".join(output["data"]["text/plain"]))
            block = ["[Code Cell]", source.strip()]
            if outputs:
                block.extend(["[Output]", "\n".join(outputs).strip()])
            if cell_buffer:
                sections.append(
                    ParsedSection(
                        title=current_title,
                        text="\n\n".join(cell_buffer),
                        section=current_title,
                        metadata={"cell_index": current_index, "content_kind": "markdown", **text_metadata},
                    )
                )
                cell_buffer = []
            sections.append(
                ParsedSection(
                    title=f"{current_title} code",
                    text="\n".join(part for part in block if part),
                    section=current_title,
                    metadata={"cell_index": current_index, "content_kind": "code", **text_metadata},
                )
            )
    if cell_buffer:
        sections.append(
            ParsedSection(
                title=current_title,
                text="\n\n".join(cell_buffer),
                section=current_title,
                metadata={"content_kind": "markdown", **text_metadata},
            )
        )
    return sections


def parse_image(path: Path) -> list[ParsedSection]:
    try:
        from paddleocr import PaddleOCR

        ocr = PaddleOCR(use_angle_cls=True, lang="ch")
        results = ocr.ocr(str(path), cls=True)
        text = "\n".join(line[1][0] for page in results for line in page)
        if text.strip():
            return [ParsedSection(title=path.stem, text=text.strip(), section=path.stem, metadata={"content_kind": "ocr"})]
    except Exception:
        pass

    try:
        import pytesseract
        from PIL import Image

        image = Image.open(path)
        text = pytesseract.image_to_string(image, lang="chi_sim+eng")
        if text.strip():
            return [ParsedSection(title=path.stem, text=text.strip(), section=path.stem, metadata={"content_kind": "ocr"})]
    except Exception as exc:
        raise RuntimeError(f"OCR dependencies unavailable for {path.name}: {exc}") from exc
    raise RuntimeError(f"No OCR text extracted from {path.name}")


def parse_with_unstructured(path: Path) -> list[ParsedSection]:
    from unstructured.partition.auto import partition

    elements = partition(filename=str(path))
    text = "\n".join(str(element) for element in elements if str(element).strip())
    return [ParsedSection(title=path.stem, text=text, section=path.stem, metadata={"content_kind": "unstructured"})] if text else []


def parse_document(path: Path) -> tuple[str, list[ParsedSection]]:
    source_type = detect_source_type(path)
    parsers = {
        "pdf": parse_pdf,
        "ppt": parse_presentation,
        "pptx": parse_presentation,
        "docx": parse_docx,
        "markdown": parse_markdown,
        "text": parse_text,
        "image": parse_image,
        "notebook": parse_notebook,
        "html": parse_html,
    }
    parser = parsers.get(source_type)
    if parser is None:
        if get_settings().enable_model_fallback:
            return source_type, clean_parsed_sections(parse_with_unstructured(path), source_type)
        raise RuntimeError(f"Unsupported source type for {path.name}: {source_type}")
    try:
        return source_type, clean_parsed_sections(parser(path), source_type)
    except Exception as exc:
        if get_settings().enable_model_fallback:
            return source_type, clean_parsed_sections(parse_with_unstructured(path), source_type)
        raise RuntimeError(f"Failed to parse {path.name} as {source_type}: {exc}") from exc


def sections_to_json(sections: list[ParsedSection]) -> list[dict[str, Any]]:
    return [asdict(section) for section in sections]


INVALID_CHAPTER_LABELS = {
    "data",
    "storage",
    "reviewmarkdown",
}


def _normalize_label(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", value.lower()).strip()


def is_date_like_label(value: str) -> bool:
    return bool(re.fullmatch(r"(?:19|20)\d{6}", value.strip()))


def is_invalid_chapter_label(value: str | None, course_name: str | None = None) -> bool:
    if not value:
        return True
    normalized = _normalize_label(value)
    if not normalized:
        return True
    if normalized.replace(" ", "") in INVALID_CHAPTER_LABELS:
        return True
    if is_date_like_label(normalized.replace(" ", "")):
        return True
    if course_name and normalized == _normalize_label(course_name):
        return True
    return False


def canonical_chapter_label(value: str, course_name: str | None = None) -> str | None:
    cleaned = re.sub(r"[_-]+", " ", value)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if not cleaned:
        return None

    match = re.search(r"\b(?:chapter|chap)\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Chapter {match.group(1)}"

    match = re.search(r"\b(?:lecture|lec|l)\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Lecture {match.group(1)}"

    match = re.search(r"\bweek\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Week {match.group(1)}"

    match = re.search(r"\blab\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Lab {match.group(1)}"

    if re.search(r"\blabs?\b|\blaboratory\b", cleaned, flags=re.IGNORECASE):
        if re.search(r"\bquestions?\b", cleaned, flags=re.IGNORECASE):
            return "Lab Questions"
        if re.search(r"\bsolutions?\b", cleaned, flags=re.IGNORECASE):
            return "Lab Solutions"
        return "Lab"

    if re.search(r"\bcourse\s*work\b|\bcoursework\b", cleaned, flags=re.IGNORECASE):
        return "Coursework"

    if re.search(r"\bz\s*table\b|\breference\b|\bformula\b|\bsummary\b|\bvisuali[sz]er\b", cleaned, flags=re.IGNORECASE):
        return "Reference"

    if re.search(r"\breview\b|\brevision\b|\brecap\b", cleaned, flags=re.IGNORECASE):
        return "Review"

    cleaned = cleaned[:80].strip()
    if cleaned and not re.search(r"[A-Za-z0-9]", cleaned):
        return "Reference"
    return None if is_invalid_chapter_label(cleaned, course_name=course_name) else cleaned


def derive_chapter(path: Path, course_name: str | None = None) -> str:
    candidates = [path.stem, path.parent.name]
    for item in candidates:
        label = canonical_chapter_label(item, course_name=course_name)
        if label and not is_invalid_chapter_label(label, course_name=course_name):
            return label
    return path.stem[:80]
